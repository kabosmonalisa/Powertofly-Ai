#!/usr/bin/env python3
"""
Build PowerToFly-GTM-Sales-Deck.pptx from decks/gtm-deck.html.

The HTML is the source of truth: this reads the real slide markup and
re-emits every slide as NATIVE PowerPoint shapes, tables and text boxes —
so Milena's team can edit every word after importing into Google Slides.

Geometry is a straight translation of the HTML's container-query units:
the slide is 100cqw wide, so 1cqw = 13.3333/100 in and 1cqw = 9.6 pt.
Colours are the ptf.css tokens. Nothing here is invented.
"""
import os, re, io, sys, urllib.request, shutil
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, '_assets')
os.makedirs(ASSETS, exist_ok=True)

# ---------------------------------------------------------------- tokens
INK       = RGBColor(0x0A, 0x0A, 0x0A)
INK2      = RGBColor(0x2A, 0x2A, 0x2A)
INK3      = RGBColor(0x5C, 0x6F, 0x69)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x0A, 0x8C, 0x66)
GREEN_MID = RGBColor(0x4F, 0xE8, 0xA9)
BG_SOFT   = RGBColor(0xFA, 0xFA, 0xF6)
LINE      = RGBColor(0xE7, 0xE7, 0xE5)   # rgba(10,10,10,.08) on white
LINE_D    = RGBColor(0x24, 0x24, 0x24)   # rgba(255,255,255,.10) on ink
D_INK2    = RGBColor(0xCF, 0xCF, 0xCF)   # 80% white on ink
D_INK3    = RGBColor(0x85, 0x85, 0x85)   # 50% white on ink
GREEN_TINT= RGBColor(0xF2, 0xF8, 0xF6)   # rgba(10,140,102,.05) on white

H_FONT = 'Inter Tight'
B_FONT = 'Inter'

CQ = 13.3333 / 100.0     # inches per cqw
PTS = 9.6                # points per cqw
PAD_X, PAD_Y = 6.5, 5.6  # the .pad safe area, in cqw
CW = 100 - 2 * PAD_X     # content width in cqw


def I(cqw):  return Inches(cqw * CQ)
def P(cqw):  return Pt(cqw * PTS)


# ---------------------------------------------------------------- assets
def asset(url_or_path):
    """Return a local file path for a remote URL or a repo-relative path."""
    if url_or_path.startswith('http'):
        name = re.sub(r'[^a-zA-Z0-9]+', '-', url_or_path.split('?')[0])[-80:] + '.jpg'
        dest = os.path.join(ASSETS, name)
        if not os.path.exists(dest):
            req = urllib.request.Request(url_or_path, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=60) as r, open(dest, 'wb') as f:
                shutil.copyfileobj(r, f)
        return _fit(dest)
    return _fit(os.path.normpath(os.path.join(HERE, url_or_path)))


def _fit(path, cap=1600):
    """Cache a downscaled JPEG so the .pptx stays small enough to import."""
    from PIL import Image
    out = os.path.join(ASSETS, '_fit-' + os.path.basename(path).rsplit('.', 1)[0] + '.jpg')
    if os.path.exists(out):
        return out
    im = Image.open(path).convert('RGB')
    if max(im.size) > cap:
        im.thumbnail((cap, cap), Image.LANCZOS)
    im.save(out, 'JPEG', quality=82, optimize=True)
    return out


def bg_url(style):
    m = re.search(r'url\(([^)]+)\)', style or '')
    return m.group(1).strip('\'"') if m else None


# ---------------------------------------------------------------- drawing
def set_spacing(run, cqw_em):
    """Letter-spacing, in 1/100 pt, via the raw rPr (python-pptx has no API)."""
    run.font._rPr.set('spc', str(int(cqw_em * 100)))


def textbox(sl, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def write(p, text, size, color, font=B_FONT, bold=False, italic=False,
          spacing=None, line=None, align=None, caps=False):
    r = p.add_run()
    r.text = text.upper() if caps else text
    r.font.size = P(size); r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    if spacing: set_spacing(r, spacing * size * PTS)
    if line: p.line_spacing = line
    if align is not None: p.alignment = align
    return r


def rect(sl, x, y, w, h, fill=None, line=None, radius=None, lw=1):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        I(x), I(y), I(w), I(h))
    if radius:
        shp.adjustments[0] = min(0.5, radius / min(w, h))
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:    shp.fill.background()
    if line: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ''
    return shp


def hline(sl, x, y, w, color=LINE, weight=1):
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), Pt(weight))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def picture(sl, x, y, w, h, src, crop_focus=0.5):
    """Cover-fit a picture into the box, cropping the overflow like CSS cover."""
    path = asset(src)
    pic = sl.shapes.add_picture(path, I(x), I(y), I(w), I(h))
    from PIL import Image  # optional; fall back to a plain stretch
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return pic
    box_r, img_r = w / h, iw / ih
    if img_r > box_r:                       # image too wide -> crop sides
        keep = box_r / img_r
        pic.crop_left = pic.crop_right = (1 - keep) / 2
    else:                                   # too tall -> crop top/bottom
        keep = img_r / box_r
        over = 1 - keep
        pic.crop_top = over * crop_focus
        pic.crop_bottom = over * (1 - crop_focus)
    return pic


# ---------------------------------------------------------------- HTML bits
def txt(el):
    return re.sub(r'\s+', ' ', el.get_text(' ', strip=True)) if el else ''


def heading_runs(el):
    """[(text, is_em)] so the italic accent survives as a real run."""
    out = []
    for node in el.children:
        if getattr(node, 'name', None) == 'em':
            out.append((txt(node), True))
        elif getattr(node, 'name', None) == 'br':
            out.append(('\n', False))
        else:
            s = re.sub(r'\s+', ' ', str(node) if isinstance(node, str) else node.get_text(' '))
            if s.strip() or s == ' ':
                out.append((s, False))
    return out


def dark(sl_el):
    return 'dark' in sl_el.get('class', [])


def C(is_dark, light, darkc):
    return darkc if is_dark else light


# ================================================================ chrome
def chrome(sl, el, d):
    """Eyebrow + headline + footer — the parts every slide shares."""
    y = PAD_Y
    eb = el.select_one('.d-eyebrow')
    if eb:
        tf = textbox(sl, PAD_X, y, CW, 1.6)
        write(para(tf, True), txt(eb), 1.15, C(d, GREEN, GREEN_MID),
              bold=True, spacing=.16, caps=True)
        y += 1.15 * 1.3 + 2.4
    h = el.select_one('.d-h1, .d-h2')
    if h:
        size = float(re.search(r'font-size:([\d.]+)cqw', h.get('style', '') or '')
                     .group(1)) if 'font-size' in (h.get('style') or '') else \
               (5.4 if 'd-h1' in h.get('class', []) else 4.2)
        mb = float(re.search(r'margin-bottom:([\d.]+)cqw', h.get('style', '') or '')
                   .group(1)) if 'margin-bottom' in (h.get('style') or '') else 2.0
        lines = max(1, len(txt(h)) / max(18, (CW * PTS) / (size * PTS * .52)))
        hh = size * 1.05 * (int(lines) + 1)
        tf = textbox(sl, PAD_X, y, CW, hh)
        p = para(tf, True)
        for t, em in heading_runs(h):
            if t == '\n':
                p = para(tf); continue
            write(p, t, size, C(d, INK, WHITE), font=H_FONT, bold=True,
                  italic=em, line=1.05, spacing=-.03)
        y += size * 1.06 * max(1, _wrapped(txt(h), size)) + mb
    ft = el.select_one('.d-foot')
    if ft:
        sp = ft.find_all('span', recursive=False)
        tf = textbox(sl, PAD_X, 100 * 9 / 16 - PAD_Y - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM)
        p = para(tf, True)
        write(p, txt(sp[0]) if sp else '', 1.15, C(d, INK3, D_INK3))
        if len(sp) > 1:
            tf2 = textbox(sl, PAD_X, 100 * 9 / 16 - PAD_Y - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM)
            write(para(tf2, True), txt(sp[-1]), 1.15, C(d, INK3, D_INK3),
                  align=PP_ALIGN.RIGHT)
    return y


def _wrapped(s, size):
    """Rough line count at this type size across the content width."""
    chars_per_line = (CW * PTS) / (size * PTS * 0.5)
    return max(1, -(-len(s) // int(chars_per_line)))


def subline(sl, el, y, d):
    p_ = el.select_one('.pad > p.d-body')
    if not p_:
        return y
    size = float(re.search(r'font-size:([\d.]+)cqw', p_.get('style', '') or '').group(1)) \
        if 'font-size' in (p_.get('style') or '') else 1.75
    mb = float(re.search(r'margin-bottom:([\d.]+)cqw', p_.get('style', '') or '').group(1)) \
        if 'margin-bottom' in (p_.get('style') or '') else 0
    h = size * 1.5 * _wrapped(txt(p_), size)
    tf = textbox(sl, PAD_X, y, CW, h)
    for t, em in heading_runs(p_):
        write(para(tf, True) if not tf.paragraphs[0].runs else tf.paragraphs[0],
              t, size, C(d, INK2, D_INK2), bold=em, line=1.5)
    return y + h + mb


# ================================================================ layouts
BOT = 100 * 9 / 16 - PAD_Y          # content bottom edge, in cqw


def L_stats(sl, el, y, d):
    items = el.select('.stat')
    cols, gap = 3, 3.0
    cw = (CW - gap * (cols - 1)) / cols
    rows = -(-len(items) // cols)
    rh = min(9.0, (BOT - 1.6 - y) / rows)
    for i, st in enumerate(items):
        cx = PAD_X + (i % cols) * (cw + gap)
        cy = y + (i // cols) * rh
        tf = textbox(sl, cx, cy, cw, 5.4)
        write(para(tf, True), txt(st.select_one('.n')), 5.2, C(d, GREEN, GREEN_MID),
              font=H_FONT, bold=True, line=1.0, spacing=-.02)
        tf2 = textbox(sl, cx, cy + 5.4, cw, 3.0)
        write(para(tf2, True), txt(st.select_one('.l')), 1.5, C(d, INK2, D_INK2), line=1.3)


def L_stats5(sl, el, y, d):
    items = el.select('.stat'); gap = 2.0
    cw = (CW - gap * 4) / 5
    for i, st in enumerate(items):
        cx = PAD_X + i * (cw + gap)
        tf = textbox(sl, cx, y, cw, 5.4)
        write(para(tf, True), txt(st.select_one('.n')), 4.4, C(d, GREEN, GREEN_MID),
              font=H_FONT, bold=True, line=1.0, spacing=-.02)
        tf2 = textbox(sl, cx, y + 4.8, cw, 3.4)
        write(para(tf2, True), txt(st.select_one('.l')), 1.4, C(d, INK2, D_INK2), line=1.3)


def L_roles(sl, el, y, d):
    cols = el.select('.rolecol'); gap = 3.0
    cw = (CW - gap * (len(cols) - 1)) / len(cols)
    for i, col in enumerate(cols):
        cx = PAD_X + i * (cw + gap)
        tf = textbox(sl, cx, y, cw, 1.5)
        write(para(tf, True), txt(col.select_one('.rlbl')), 1.0, C(d, GREEN, GREEN_MID),
              bold=True, spacing=.14, caps=True)
        yy = y + 2.2
        for li in col.select('li'):
            hline(sl, cx, yy, cw, C(d, LINE, LINE_D))
            tf = textbox(sl, cx, yy + .55, cw, 1.7)
            write(para(tf, True), txt(li), 1.2, C(d, INK2, D_INK2), line=1.25)
            yy += 2.44


def L_facerow(sl, el, d):
    fr = el.select_one('.facerow')
    if not fr: return
    imgs = fr.select('img'); sz = 5.0; gap = 1.1
    yy = BOT - 1.9 - sz
    for i, im in enumerate(imgs):
        x = PAD_X + i * (sz + gap)
        pic = picture(sl, x, yy, sz, sz, im['src'])
        _make_oval(pic)
    lbl = fr.select_one('.fr-label')
    if lbl:
        lx = PAD_X + len(imgs) * (sz + gap) + .6
        tf = textbox(sl, lx, yy, max(4, 100 - PAD_X - lx), sz, MSO_ANCHOR.MIDDLE)
        write(para(tf, True), txt(lbl), 1.2, INK3, line=1.3)


def _make_oval(pic):
    """Crop a picture to a circle by REPLACING its geometry.

    Appending a second <a:prstGeom> leaves two geometries in one <a:spPr>.
    PowerPoint forgives that; Google Slides rejects the whole file."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    spPr = pic._element.spPr
    new = parse_xml('<a:prstGeom %s prst="ellipse"><a:avLst/></a:prstGeom>' % nsdecls('a'))
    for tag in ('a:prstGeom', 'a:custGeom'):
        old = spPr.find(qn(tag))
        if old is not None:
            spPr.replace(old, new)
            return
    xfrm = spPr.find(qn('a:xfrm'))
    spPr.insert(list(spPr).index(xfrm) + 1 if xfrm is not None else 0, new)


def L_probs(sl, el, y, d):
    cards = el.select('.pcard'); gap_x, gap_y = 2.6, 2.0
    cw = (CW - gap_x) / 2
    ch = (BOT - 1.6 - y - gap_y) / 2
    for i, c in enumerate(cards):
        cx = PAD_X + (i % 2) * (cw + gap_x)
        cy = y + (i // 2) * (ch + gap_y)
        hline(sl, cx, cy, cw, C(d, LINE, LINE_D), 2)
        tf = textbox(sl, cx, cy + 1.6, cw, 3.0)
        write(para(tf, True), txt(c.select_one('.ps')), 2.6, C(d, GREEN, GREEN_MID),
              font=H_FONT, bold=True, line=1.0)
        tf = textbox(sl, cx, cy + 4.4, cw, 2.4)
        write(para(tf, True), txt(c.select_one('h3')), 1.8, C(d, INK, WHITE),
              font=H_FONT, bold=True)
        tf = textbox(sl, cx, cy + 6.8, cw, ch - 6.8)
        write(para(tf, True), txt(c.select_one('p')), 1.3, C(d, INK2, D_INK2), line=1.4)


def L_svc(sl, el, y, d):
    cols = el.select('.svccol'); gap = 4.0
    cw = (CW - gap) / 2
    for i, col in enumerate(cols):
        cx = PAD_X + i * (cw + gap); yy = y
        write(para(textbox(sl, cx, yy, cw, 1.6), True), txt(col.select_one('.svcnum')),
              1.4, GREEN, font=H_FONT, bold=True); yy += 2.0
        write(para(textbox(sl, cx, yy, cw, 3.0), True), txt(col.select_one('h3')),
              2.4, INK, font=H_FONT, bold=True, spacing=-.01); yy += 3.4
        p0 = col.find('p', recursive=False)
        if p0:
            h = 1.4 * 1.45 * _wrapped(txt(p0), 1.4)
            write(para(textbox(sl, cx, yy, cw, h), True), txt(p0), 1.4, INK2, line=1.45)
            yy += h + 1.2
        for child in col.find_all(['ul', 'div'], recursive=False):
            if 'svclbl' in (child.get('class') or []):
                write(para(textbox(sl, cx, yy, cw, 1.4), True), txt(child), 1.05, INK3,
                      bold=True, spacing=.14, caps=True); yy += 2.0
            elif child.name == 'ul':
                for li in child.select('li'):
                    tf = textbox(sl, cx, yy, cw, 2.0)
                    p = para(tf, True)
                    write(p, '✓  ', 1.35, GREEN, bold=True)
                    write(p, txt(li), 1.35, INK, line=1.35)
                    yy += 1.35 * 1.35 + .7


def L_proc(sl, el, y, d, cols=None):
    steps = el.select('.pstep'); cols = cols or len(steps); gap = 2.4
    cw = (CW - gap * (cols - 1)) / cols
    for i, st in enumerate(steps):
        cx = PAD_X + (i % cols) * (cw + gap)
        cy = y + (i // cols) * 11
        write(para(textbox(sl, cx, cy, cw, 3.2), True), txt(st.select_one('.pn')),
              2.8, C(d, GREEN, GREEN_MID), font=H_FONT, bold=True, line=1.0)
        write(para(textbox(sl, cx, cy + 3.5, cw, 4.4), True), txt(st.select_one('h3')),
              1.7 if cols < 5 else 1.5, C(d, INK, WHITE), font=H_FONT, bold=True, line=1.15)
        write(para(textbox(sl, cx, cy + 6.4, cw, 8), True), txt(st.select_one('p')),
              1.25, C(d, INK2, D_INK2), line=1.4)
    band = el.select_one('.photoband')
    if band:
        top = y + 12.5
        picture(sl, PAD_X, top, CW, BOT - 1.9 - top, bg_url(band.get('style')), .35)


def L_bens(sl, el, y, d):
    items = el.select('.ben'); gap = 2.4
    cw = (CW - gap * 3) / 4
    for i, b in enumerate(items):
        cx = PAD_X + i * (cw + gap)
        c = rect(sl, cx, y, 3.4, 3.4, BG_SOFT, LINE, radius=1.7)
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(para(tf, True), txt(b.select_one('.bnum')), 1.5, GREEN,
              font=H_FONT, bold=True, align=PP_ALIGN.CENTER)
        write(para(textbox(sl, cx, y + 4.4, cw, 4.0), True), txt(b.select_one('h3')),
              1.7, INK, font=H_FONT, bold=True, line=1.15)
        write(para(textbox(sl, cx, y + 8.0, cw, BOT - 1.9 - y - 8.0), True),
              txt(b.select_one('p')), 1.25, INK2, line=1.4)


def L_grid6(sl, el, y, d):
    cards = el.select('.dcard'); dense = 'dense' in (el.select_one('.grid6').get('class'))
    gx, gy = 2.6, 1.4
    cw = (CW - gx) / 2
    rows = -(-len(cards) // 2)
    ch = (BOT - 1.6 - y - gy * (rows - 1)) / rows
    for i, c in enumerate(cards):
        cx = PAD_X + (i % 2) * (cw + gx)
        cy = y + (i // 2) * (ch + gy)
        rect(sl, cx, cy, cw, ch, None, LINE, radius=1.2)
        rect(sl, cx, cy, .5, ch, GREEN)
        tf = textbox(sl, cx + 1.8, cy + 1.2, cw - 3.4, ch - 2.0, MSO_ANCHOR.MIDDLE)
        write(para(tf, True), txt(c.select_one('.dn')), 1.05 if dense else 1.2, GREEN,
              font=H_FONT, bold=True)
        write(para(tf), txt(c.select_one('h3')), 1.5 if dense else 1.7, INK,
              font=H_FONT, bold=True, line=1.1)
        pp = c.select_one('p')
        if pp:
            write(para(tf), txt(pp), 1.15 if dense else 1.25, INK2, line=1.3)


def L_cmp(sl, el, y, d):
    cols = el.select('.cmpcol'); gap = 1.6
    cw = (CW - gap * 3) / 4
    ch = BOT - 1.6 - y
    for i, c in enumerate(cols):
        cx = PAD_X + i * (cw + gap)
        us = 'cmpcol--us' in c.get('class', [])
        rect(sl, cx, y, cw, ch, GREEN_TINT if us else None, GREEN if us else LINE, radius=1.2)
        tf = textbox(sl, cx + 1.6, y + 1.8, cw - 3.2, ch - 3.0)
        write(para(tf, True), txt(c.select_one('.cat')), .95, INK3, bold=True,
              spacing=.12, caps=True)
        write(para(tf), txt(c.select_one('h3')), 1.7, INK, font=H_FONT, bold=True, line=1.1)
        for li in c.select('li'):
            p = para(tf)
            yes = 'y' in li.get('class', [])
            write(p, '✓  ' if yes else '✕  ', 1.15, GREEN if yes else INK3, bold=True)
            write(p, txt(li), 1.15, INK2, line=1.3)


def L_divider(sl, el, y, d):
    ph = el.select_one('.slide-photo')
    if ph:
        picture(sl, 40, 0, 60, 100 * 9 / 16, bg_url(ph.get('style')), .55)
        grad = rect(sl, 40, 0, 60, 100 * 9 / 16, INK)
        _gradient(grad)
    # eyebrow / headline / sub are drawn on top
    eb = el.select_one('.d-eyebrow')
    write(para(textbox(sl, PAD_X, PAD_Y, CW, 1.6), True), txt(eb), 1.15, GREEN_MID,
          bold=True, spacing=.16, caps=True)
    h = el.select_one('.d-h2')
    size = float(re.search(r'font-size:([\d.]+)cqw', h.get('style', '')).group(1))
    tf = textbox(sl, PAD_X, 100 * 9 / 16 / 2 - size * .8, 38, size * 1.4)
    p = para(tf, True)
    for t, em in heading_runs(h):
        write(p, t, size, WHITE, font=H_FONT, bold=True, italic=em, line=1.03, spacing=-.03)
    sub = el.select_one('.div-sub')
    if sub:
        write(para(textbox(sl, PAD_X, 100 * 9 / 16 / 2 + size * .75, 32, 6), True),
              txt(sub), 1.7, D_INK2, line=1.35)
    ft = el.select_one('.d-foot')
    if ft:
        sp = ft.find_all('span', recursive=False)
        write(para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True),
              txt(sp[0]), 1.15, D_INK3)
        write(para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True),
              txt(sp[-1]), 1.15, D_INK3, align=PP_ALIGN.RIGHT)


def _gradient(shape):
    """Ink → transparent, left to right — the divider scrim."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill'):
        e = spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    xml = ('<a:gradFill %s rotWithShape="1">'
           '<a:gsLst>'
           '<a:gs pos="0"><a:srgbClr val="0A0A0A"><a:alpha val="100000"/></a:srgbClr></a:gs>'
           '<a:gs pos="22000"><a:srgbClr val="0A0A0A"><a:alpha val="72000"/></a:srgbClr></a:gs>'
           '<a:gs pos="60000"><a:srgbClr val="0A0A0A"><a:alpha val="28000"/></a:srgbClr></a:gs>'
           '<a:gs pos="100000"><a:srgbClr val="0A0A0A"><a:alpha val="12000"/></a:srgbClr></a:gs>'
           '</a:gsLst><a:lin ang="0" scaled="0"/></a:gradFill>') % nsdecls('a')
    ln = spPr.find(qn('a:ln'))
    (spPr.insert(list(spPr).index(ln), parse_xml(xml)) if ln is not None
     else spPr.append(parse_xml(xml)))


def L_deckmock(sl, el, x, w):
    """The product mock — dark card with rows."""
    mock = el.select_one('.deckmock')
    if not mock: return
    rows = mock.select('.dm-row')
    h = 1.9 * 2 + 4.2 + len(rows) * (3.4 + .75)
    y = (100 * 9 / 16 - h) / 2
    rect(sl, x, y, w, h, RGBColor(0x0C, 0x14, 0x10), RGBColor(0x1A, 0x24, 0x20), radius=1.8)
    ic = rect(sl, x + 1.9, y + 1.9, 2.9, 2.9, GREEN_MID, radius=.85)
    tf = textbox(sl, x + 1.9 + 3.9, y + 1.9, w - 12, 1.6)
    write(para(tf, True), txt(mock.select_one('.dm-ttl')), 1.2, WHITE, font=H_FONT, bold=True)
    write(para(textbox(sl, x + 1.9 + 3.9, y + 3.4, w - 12, 1.4), True),
          txt(mock.select_one('.dm-sub')), 1.0, RGBColor(0x66, 0x66, 0x66))
    write(para(textbox(sl, x + 1.9, y + 1.9, w - 3.8, 1.6), True),
          txt(mock.select_one('.dm-live')), .9, RGBColor(0x80, 0x80, 0x80),
          spacing=.1, align=PP_ALIGN.RIGHT)
    ry = y + 1.9 + 4.2
    for r in rows:
        active = 'active' in r.get('class', [])
        rect(sl, x + 1.9, ry, w - 3.8, 3.4,
             RGBColor(0x11, 0x2C, 0x22) if active else RGBColor(0x14, 0x1B, 0x18),
             GREEN_MID if active else None, radius=1.1)
        sp = r.find_all('span', recursive=False)
        tf = textbox(sl, x + 3.2, ry, w - 6.4, 3.4, MSO_ANCHOR.MIDDLE)
        write(para(tf, True), txt(sp[0]), 1.15, WHITE if active else D_INK2)
        tag = r.select_one('.tag')
        if tag:
            tf2 = textbox(sl, x + 3.2, ry, w - 6.4, 3.4, MSO_ANCHOR.MIDDLE)
            write(para(tf2, True), txt(tag), .92, GREEN_MID if active else D_INK3,
                  bold=active, align=PP_ALIGN.RIGHT)
        ry += 3.4 + .75


def L_cover(sl, el, d):
    lg = el.select_one('.d-logo')
    tf = textbox(sl, PAD_X, PAD_Y, 30, 2.4)
    p = para(tf, True)
    write(p, 'PowerToFly ', 1.7, WHITE, font=H_FONT, bold=True, spacing=-.02)
    write(p, 'AI', 1.7, GREEN_MID, font=H_FONT, bold=True, spacing=-.02)
    left = el.select_one('.covergrid > div')
    colw = (CW - 4) * (1.25 / 2.0)
    y = 100 * 9 / 16 * .33
    write(para(textbox(sl, PAD_X, y, colw, 1.6), True), txt(left.select_one('.d-eyebrow')),
          1.15, GREEN_MID, bold=True, spacing=.16, caps=True)
    h = left.select_one('.d-h1'); size = 4.6
    tfh = textbox(sl, PAD_X, y + 3.0, colw, size * 3.4)
    p = para(tfh, True)
    for t, em in heading_runs(h):
        write(p, t, size, WHITE, font=H_FONT, bold=True, italic=em, line=1.03, spacing=-.03)
    body = left.select_one('.d-body')
    write(para(textbox(sl, PAD_X, y + 3.0 + size * 3.2, colw, 8), True), txt(body),
          1.75, D_INK2, line=1.5)
    L_deckmock(sl, el, PAD_X + colw + 4, CW - colw - 4)
    ft = el.select_one('.d-foot'); sp = ft.find_all('span', recursive=False)
    write(para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True),
          txt(sp[0]), 1.15, D_INK3)
    write(para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True),
          txt(sp[-1]), 1.15, D_INK3, align=PP_ALIGN.RIGHT)


def L_reach(sl, el, y, d):
    svg = el.select_one('.reachmap svg')
    mw = CW * (1.55 / 2.45) - 1.8
    mh = mw / 2.0
    my = y + 2
    for c in svg.select('circle'):
        cx = float(c['cx']) / 200 * mw + PAD_X
        cy = float(c['cy']) / 100 * mh + my
        r = float(c['r']) / 200 * mw
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, I(cx - r), I(cy - r), I(r * 2), I(r * 2))
        lit = 'green' in (c.get('fill') or '')
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN_MID if lit else RGBColor(0x3A, 0x3A, 0x3A)
        dot.line.fill.background(); dot.shadow.inherit = False
    sx = PAD_X + mw + 3.5
    sw = (CW - mw - 3.5 - 2.6) / 2
    for i, st in enumerate(el.select('.rstat')):
        cx = sx + (i % 2) * (sw + 2.6)
        cy = my + (i // 2) * 8.0
        write(para(textbox(sl, cx, cy, sw, 3.6), True), txt(st.select_one('.rn')),
              3.4, GREEN_MID, font=H_FONT, bold=True, line=1.0)
        write(para(textbox(sl, cx, cy + 4.0, sw, 4.0), True), txt(st.select_one('.rl')),
              1.15, D_INK2, line=1.3)


def L_talent(sl, el, d):
    colw = (CW - 4) * (1.05 / 2.0)
    y = PAD_Y
    write(para(textbox(sl, PAD_X, y, colw, 1.6), True), txt(el.select_one('.d-eyebrow')),
          1.15, GREEN_MID, bold=True, spacing=.16, caps=True)
    h = el.select_one('.d-h2'); size = 2.7
    tf = textbox(sl, PAD_X, y + 3.0, colw, size * 3.4)
    p = para(tf, True)
    for t, em in heading_runs(h):
        write(p, t, size, WHITE, font=H_FONT, bold=True, italic=em, line=1.05, spacing=-.03)
    yy = y + 3.0 + size * 2.3 + 1.2
    for lay in el.select('.layer'):
        rect(sl, PAD_X, yy, .45, 5.6, GREEN)
        write(para(textbox(sl, PAD_X + 2.2, yy, colw - 2.2, 2.6), True),
              txt(lay.select_one('h3')), 2.1, WHITE, font=H_FONT, bold=True)
        write(para(textbox(sl, PAD_X + 2.2, yy + 2.7, colw - 2.2, 3.6), True),
              txt(lay.select_one('p')), 1.3, D_INK2, line=1.38)
        yy += 5.6 + 1.3
    L_deckmock(sl, el, PAD_X + colw + 4, CW - colw - 4)


def L_audience(sl, el, y, d):
    colw = CW * (.85 / 2.0) - 2
    big = el.select_one('.aud-big')
    write(para(textbox(sl, PAD_X, y, colw, 5.4), True), txt(big.select_one('.n')),
          5.0, GREEN, font=H_FONT, bold=True, line=1.0)
    write(para(textbox(sl, PAD_X, y + 5.8, colw, 8), True), txt(big.select_one('.l')),
          1.4, INK2, line=1.4)
    ph = el.select_one('.aud-photo')
    if ph:
        picture(sl, PAD_X, y + 14.2, colw, BOT - 2.0 - (y + 14.2), bg_url(ph.get('style')))
    rx = PAD_X + colw + 4; rw = CW - colw - 4
    yy = y
    for li in el.select('.checks li'):
        b = li.find('b')
        tf = textbox(sl, rx, yy, rw, 6)
        p = para(tf, True)
        write(p, '✓  ', 1.35, GREEN, bold=True)
        if b:
            write(p, txt(b) + ' ', 1.35, INK, bold=True, line=1.4)
            rest = txt(li).replace(txt(b), '', 1).lstrip(' —')
            write(p, '— ' + rest, 1.35, INK2, line=1.4)
        else:
            write(p, txt(li), 1.35, INK2, line=1.4)
        yy += 1.35 * 1.4 * (1 + _wrapped(txt(li), 1.35) - 1) + 1.5


def L_cal(sl, el, y, d):
    cards = el.select('.calcard'); gx = gy = 1.8
    cw = (CW - gx * 2) / 3
    ch = (BOT - 1.6 - y - gy) / 2
    for i, c in enumerate(cards):
        cx = PAD_X + (i % 3) * (cw + gx); cy = y + (i // 3) * (ch + gy)
        rect(sl, cx, cy, cw, ch, None, LINE, radius=1.2)
        tf = textbox(sl, cx + 1.7, cy + 1.4, cw - 3.4, ch - 2.4, MSO_ANCHOR.MIDDLE)
        write(para(tf, True), txt(c.select_one('.cmonth')), 1.0, GREEN, bold=True,
              spacing=.12, caps=True)
        write(para(tf), txt(c.select_one('h3')), 1.7, INK, font=H_FONT, bold=True, line=1.1)
        write(para(tf), txt(c.select_one('p')), 1.2, INK3, line=1.25)


def L_dedicated(sl, el, d):
    colw = CW * (1.05 / 2.0) - 2
    y = PAD_Y
    write(para(textbox(sl, PAD_X, y, colw, 1.6), True), txt(el.select_one('.d-eyebrow')),
          1.15, GREEN, bold=True, spacing=.16, caps=True)
    h = el.select_one('.d-h2')
    tf = textbox(sl, PAD_X, y + 3.0, colw, 6)
    p = para(tf, True)
    for t, em in heading_runs(h):
        write(p, t, 3.0, INK, font=H_FONT, bold=True, italic=em, line=1.05, spacing=-.03)
    body = el.select_one('.d-body')
    write(para(textbox(sl, PAD_X, y + 7.4, colw, 6), True), txt(body), 1.5, INK2, line=1.5)
    yy = y + 13.0
    for row in el.select('.solnrow'):
        hline(sl, PAD_X, yy, colw, LINE, 2)
        write(para(textbox(sl, PAD_X, yy + .85, 5, 2.6), True), txt(row.find('b')),
              1.9, GREEN, font=H_FONT, bold=True)
        write(para(textbox(sl, PAD_X + 6.6, yy + .95, colw - 6.6, 5), True),
              txt(row.find('span')), 1.35, INK2, line=1.35)
        yy += 5.4
    vis = el.select_one('.visual')
    if vis and bg_url(vis.get('style')):
        side = min(CW - colw - 4, BOT - PAD_Y - 2)
        picture(sl, 100 - PAD_X - side, (100 * 9 / 16 - side) / 2, side, side,
                bg_url(vis.get('style')))


def L_media(sl, el, y, d):
    cards = el.select('.mediacard'); gap = 3.0
    cw = (CW - gap) / 2
    th = cw * 9 / 16
    for i, c in enumerate(cards):
        cx = PAD_X + i * (cw + gap)
        thumb = c.select_one('.media-thumb')
        u = bg_url(thumb.get('style'))
        if u: picture(sl, cx, y, cw, th, u)
        else: rect(sl, cx, y, cw, th, INK, radius=1.6)
        pl = sl.shapes.add_shape(MSO_SHAPE.OVAL, I(cx + cw / 2 - 2.5), I(y + th / 2 - 2.5),
                                 I(5), I(5))
        pl.fill.solid(); pl.fill.fore_color.rgb = WHITE
        pl.line.fill.background(); pl.shadow.inherit = False
        tri = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, I(cx + cw / 2 - .8),
                                  I(y + th / 2 - 1.0), I(1.8), I(2.0))
        tri.rotation = 90
        tri.fill.solid(); tri.fill.fore_color.rgb = INK
        tri.line.fill.background(); tri.shadow.inherit = False
        write(para(textbox(sl, cx, y + th + 1.3, cw, 3), True), txt(c.select_one('h3')),
              2.0, INK, font=H_FONT, bold=True, spacing=-.01)
        write(para(textbox(sl, cx, y + th + 4.2, cw, 6), True), txt(c.select_one('p')),
              1.35, INK2, line=1.42)


def L_cases(sl, el, y, d):
    cards = el.select('.case'); gap = 2.6
    cw = (CW - gap * 2) / 3
    ch = BOT - 1.6 - y
    for i, c in enumerate(cards):
        cx = PAD_X + i * (cw + gap)
        rect(sl, cx, y, cw, ch, None, LINE, radius=1.4)
        tf = textbox(sl, cx + 2.2, y + 2.2, cw - 4.4, 6)
        write(para(tf, True), txt(c.select_one('h3')), 2.0, INK, font=H_FONT, bold=True)
        write(para(tf), txt(c.select_one('.tags')), 1.1, INK3, line=1.3)
        yy = y + 8.4
        for m in c.select('.cmetric'):
            write(para(textbox(sl, cx + 2.2, yy, 6, 2.6), True), txt(m.select_one('.n')),
                  2.1, GREEN, font=H_FONT, bold=True)
            write(para(textbox(sl, cx + 2.2 + 6.4, yy + .5, cw - 8.8, 3), True),
                  txt(m.select_one('.l')), 1.2, INK2, line=1.25)
            yy += 3.4


def L_tgrid(sl, el, y, d):
    cards = el.select('.tqcard'); gx, gy = 2.6, 2.0
    cw = (CW - gx) / 2
    ch = (BOT - 1.6 - y - gy) / 2
    for i, c in enumerate(cards):
        cx = PAD_X + (i % 2) * (cw + gx); cy = y + (i // 2) * (ch + gy)
        rect(sl, cx, cy, cw, ch, None, LINE, radius=1.4)
        tf = textbox(sl, cx + 2.2, cy + 2.0, cw - 4.4, ch - 4.0)
        write(para(tf, True), txt(c.select_one('.stars')), 1.2, GREEN, spacing=.12)
        write(para(tf), txt(c.select_one('blockquote')), 1.35, INK, line=1.45)
        write(para(tf), txt(c.select_one('.meta')), 1.1, INK3, line=1.3)


def L_logos(sl, el, y, d):
    items = el.select('.lg'); gap = 2.4
    cw = (CW - gap * 3) / 4
    ch = 7.0
    for i, lg in enumerate(items):
        cx = PAD_X + (i % 4) * (cw + gap); cy = y + 1.5 + (i // 4) * (ch + gap)
        s = rect(sl, cx, cy, cw, ch, BG_SOFT, LINE, radius=1.2)
        tf = s.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(para(tf, True), txt(lg), 1.3, INK3, bold=True, align=PP_ALIGN.CENTER)


def L_pricing(sl, el, y, d):
    cols = el.select('.pcol'); gap = 2.4
    cw = (CW - gap * 2) / 3
    ch = BOT - 1.6 - y
    for i, c in enumerate(cols):
        cx = PAD_X + i * (cw + gap)
        rect(sl, cx, y, cw, ch, None, LINE, radius=1.4)
        tf = textbox(sl, cx + 2.0, y + 2.2, cw - 4.0, ch - 4.0)
        write(para(tf, True), txt(c.select_one('.peyebrow')), 1.0, GREEN, bold=True,
              spacing=.12, caps=True)
        write(para(tf), txt(c.select_one('h3')), 1.9, INK, font=H_FONT, bold=True, line=1.1)
        lead = c.find('p', recursive=False)
        if lead:
            p = para(tf)
            for node in lead.children:
                nm = getattr(node, 'name', None)
                s = txt(node) if nm else re.sub(r'\s+', ' ', str(node))
                if not s.strip(): continue
                if nm == 'b':
                    write(p, s, 1.9, GREEN, font=H_FONT, bold=True)
                else:
                    write(p, s, 1.15, INK3, line=1.35)
        for row in c.select('.prow'):
            sp = row.find_all(['span', 'b'], recursive=False)
            p = para(tf)
            write(p, txt(sp[0]) + '   ', 1.2, INK2, line=1.4)
            if len(sp) > 1:
                write(p, txt(sp[1]), 1.2, GREEN, font=H_FONT, bold=True)
        ul = c.select_one('ul.checks')
        if ul:
            for li in ul.select('li'):
                p = para(tf)
                write(p, '✓  ', 1.1, GREEN, bold=True)
                write(p, txt(li), 1.1, INK, line=1.35)
        fo = c.select_one('.pfoot')
        if fo:
            write(para(tf), txt(fo), 1.0, INK3, line=1.3)


def L_dtbl(sl, el, y, d):
    tbl_el = el.select_one('.dtbl')
    head = tbl_el.select('thead th')
    body = tbl_el.select('tbody tr')
    ncol, nrow = len(head), len(body) + 1
    gh = 100 * 9 / 16
    h = min(BOT - 1.6 - y, 3.2 + len(body) * 2.6)
    gt = sl.shapes.add_table(nrow, ncol, I(PAD_X), I(y), I(CW), I(h)).table
    widths = []
    for th in head:
        m = re.search(r'width:\s*([\d.]+)%', th.get('style', '') or '')
        widths.append(float(m.group(1)) if m else None)
    known = sum(w for w in widths if w); rest = [i for i, w in enumerate(widths) if not w]
    for i, w in enumerate(widths):
        pct = w if w else (100 - known) / max(1, len(rest))
        gt.columns[i].width = I(CW * pct / 100)
    gt.rows[0].height = I(3.2)
    for r in range(1, nrow):
        gt.rows[r].height = I((h - 3.2) / max(1, len(body)))
    us_cols = {i for i, th in enumerate(head) if 'us' in th.get('class', [])}
    for ci, th in enumerate(head):
        cell = gt.cell(0, ci)
        _cell_fill(cell, WHITE)
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = I(.5); tf.margin_top = tf.margin_bottom = I(.2)
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        b = th.find('b')
        p = para(tf, True)
        if b:
            write(p, txt(b), 1.35, GREEN if ci in us_cols else INK, font=H_FONT,
                  bold=True, spacing=-.01)
            rest_txt = txt(th).replace(txt(b), '', 1).strip()
            if rest_txt:
                write(para(tf), rest_txt, .92, INK3, bold=True, spacing=.1, caps=True)
        else:
            write(p, txt(th), .92, INK3, bold=True, spacing=.1, caps=True)
    for ri, tr in enumerate(body, start=1):
        for ci, td in enumerate(tr.find_all('td')):
            cell = gt.cell(ri, ci)
            cls = td.get('class', [])
            _cell_fill(cell, GREEN_TINT if ('us' in cls or ci in us_cols) else WHITE)
            tf = cell.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = I(.5); tf.margin_top = tf.margin_bottom = I(.25)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            col = INK if ci == 0 else (GREEN if 'y' in cls else (INK3 if 'n' in cls else INK2))
            first = True
            for chunk in td.get_text('\n', strip=True).split('\n'):
                p = para(tf, first); first = False
                write(p, chunk, 1.05, col, bold=(ci == 0 or 'y' in cls), line=1.3)
    _strip_table_style(gt)


def _cell_fill(cell, rgb):
    cell.fill.solid(); cell.fill.fore_color.rgb = rgb


def _strip_table_style(tbl):
    """Kill the banded-row theme so the deck's own hairlines read."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '1'); tblPr.set('bandRow', '0')
    for e in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(e)
    for tr in tbl._tbl.findall(qn('a:tr')):
        for tc in tr.findall(qn('a:tc')):
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None: continue
            # CT_TableCellProperties is a sequence: lnL, lnR, lnT, lnB, ... then fill.
            # Appending after the fill produces a file Google Slides refuses to open.
            tcPr.insert(0, parse_xml(
                '<a:lnB %s w="9525" cap="flat"><a:solidFill><a:srgbClr val="E7E7E5"/>'
                '</a:solidFill></a:lnB>' % nsdecls('a')))


def L_contacts(sl, el, y, d):
    cs = el.select('.contact'); gap = 4.0
    cw = (CW * .8 - gap) / 2
    for i, c in enumerate(cs):
        cx = PAD_X + i * (cw + gap)
        o = sl.shapes.add_shape(MSO_SHAPE.OVAL, I(cx), I(y + 1), I(5.5), I(5.5))
        o.fill.solid(); o.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        o.line.color.rgb = RGBColor(0x40, 0x40, 0x40); o.line.width = Pt(1)
        o.shadow.inherit = False
        tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(para(tf, True), txt(c.select_one('.ph')), 1.9, D_INK3, font=H_FONT,
              bold=True, align=PP_ALIGN.CENTER)
        tx = cx + 7.3
        write(para(textbox(sl, tx, y + 1.2, cw - 7.3, 3), True), txt(c.select_one('h3')),
              2.3, WHITE, font=H_FONT, bold=True)
        write(para(textbox(sl, tx, y + 4.2, cw - 7.3, 3), True), txt(c.select_one('.role')),
              1.4, D_INK2, line=1.3)
        write(para(textbox(sl, tx, y + 6.2, cw - 7.3, 3), True), txt(c.select_one('.email')),
              1.4, GREEN_MID)


def L_cta(sl, el, d):
    cc = el.select_one('.cta-center')
    mid = 100 * 9 / 16 / 2
    write(para(textbox(sl, PAD_X, mid - 8.0, CW, 2), True), txt(cc.select_one('.d-eyebrow')),
          1.15, GREEN_MID, bold=True, spacing=.16, caps=True, align=PP_ALIGN.CENTER)
    h = cc.select_one('.d-h2')
    tf = textbox(sl, PAD_X, mid - 5.4, CW, 11)
    p = para(tf, True); p.alignment = PP_ALIGN.CENTER
    for t, em in heading_runs(h):
        if t == '\n':
            p = para(tf); p.alignment = PP_ALIGN.CENTER; continue
        write(p, t, 4.0, WHITE, font=H_FONT, bold=True, italic=em, line=1.05,
              spacing=-.03, align=PP_ALIGN.CENTER)
    btn = cc.select_one('.btn-pill')
    if btn:
        bw = 26
        s = rect(sl, 50 - bw / 2, mid + 6.6, bw, 4.8, WHITE, radius=2.4)
        tfb = s.text_frame; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(para(tfb, True), txt(btn), 1.7, INK, bold=True, align=PP_ALIGN.CENTER)
    ft = el.select_one('.d-foot'); sp = ft.find_all('span', recursive=False)
    p = para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True)
    write(p, 'PowerToFly ', 1.4, WHITE, font=H_FONT, bold=True)
    write(p, 'AI', 1.4, GREEN_MID, font=H_FONT, bold=True)
    write(para(textbox(sl, PAD_X, BOT - 1.4, CW, 1.4, MSO_ANCHOR.BOTTOM), True),
          txt(sp[-1]), 1.15, D_INK3, align=PP_ALIGN.RIGHT)


def L_twolist(sl, el, y, d):
    cols = el.select('.tl'); gap = 4.0
    cw = (CW - gap) / 2
    for i, c in enumerate(cols):
        cx = PAD_X + i * (cw + gap)
        write(para(textbox(sl, cx, y, cw, 2.6), True), txt(c.select_one('h3')),
              1.9, INK, font=H_FONT, bold=True, spacing=-.01)
        sub = c.select_one('.tl-sub')
        yy = y + 2.8
        if sub:
            write(para(textbox(sl, cx, yy, cw, 3), True), txt(sub), 1.15, INK3, line=1.35)
            yy += 2.6
        for li in c.select('li'):
            hline(sl, cx, yy, cw, LINE)
            n = _wrapped(txt(li), 1.25)
            write(para(textbox(sl, cx, yy + .7, cw, 1.25 * 1.35 * n + .6), True),
                  txt(li), 1.25, INK2, line=1.35)
            yy += 1.25 * 1.35 * n + 1.4


def L_persona(sl, el, y, d):
    pr = el.select_one('.persona')
    lw = CW * (.78 / 2.0) - 1.8
    rx = PAD_X + lw + 3.6; rw = CW - lw - 3.6
    idc = pr.select_one('.p-id')
    o = sl.shapes.add_shape(MSO_SHAPE.OVAL, I(PAD_X), I(y), I(6), I(6))
    o.fill.solid(); o.fill.fore_color.rgb = BG_SOFT
    o.line.color.rgb = LINE; o.line.width = Pt(1); o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(para(tf, True), txt(idc.select_one('.p-mono')), 2.4, GREEN, font=H_FONT,
          bold=True, align=PP_ALIGN.CENTER)
    yy = y + 7.2
    write(para(textbox(sl, PAD_X, yy, lw, 3.6), True), txt(idc.select_one('h3')),
          2.8, INK, font=H_FONT, bold=True, spacing=-.02); yy += 3.6
    write(para(textbox(sl, PAD_X, yy, lw, 4), True), txt(idc.select_one('.p-role')),
          1.3, INK2, line=1.35); yy += 1.3 * 1.35 * _wrapped(txt(idc.select_one('.p-role')), 1.3) + 1.4
    for meta in idc.select('.p-meta'):
        lbl = meta.select_one('.pm-l')
        write(para(textbox(sl, PAD_X, yy, lw, 1.4), True), txt(lbl), .9, INK3,
              bold=True, spacing=.12, caps=True)
        rest = txt(meta).replace(txt(lbl), '', 1).strip()
        n = _wrapped(rest, 1.1)
        write(para(textbox(sl, PAD_X, yy + 1.5, lw, 1.1 * 1.4 * n + 1), True),
              rest, 1.1, INK2, line=1.4)
        yy += 1.5 + 1.1 * 1.4 * n + 1.0
    q = pr.select_one('.p-quote')
    qn_ = _wrapped(txt(q), 1.7)
    rect(sl, rx, y, .4, 1.7 * 1.32 * qn_ + .6, GREEN)
    write(para(textbox(sl, rx + 1.8, y, rw - 1.8, 1.7 * 1.4 * qn_ + 1), True), txt(q),
          1.7, INK, font=H_FONT, bold=True, line=1.32, spacing=-.01)
    ly = y + 1.7 * 1.34 * qn_ + 2.2
    lists = pr.select('.p-lists > div')
    lcw = (rw - 2.6) / 2
    for i, col in enumerate(lists):
        cx = rx + i * (lcw + 2.6)
        write(para(textbox(sl, cx, ly, lcw, 1.5), True), txt(col.select_one('.rlbl')),
              .95, GREEN, bold=True, spacing=.13, caps=True)
        yy2 = ly + 2.0
        for li in col.select('li'):
            n = _wrapped(txt(li), 1.1)
            tf = textbox(sl, cx, yy2, lcw, 1.1 * 1.32 * n + .8)
            p = para(tf, True)
            write(p, '—  ', 1.1, GREEN)
            write(p, txt(li), 1.1, INK2, line=1.32)
            yy2 += 1.1 * 1.32 * n + .9
    pitch = pr.select_one('.p-pitch')
    if pitch:
        py = BOT - 1.9 - 4.6
        hline(sl, rx, py, rw, LINE, 2)
        tf = textbox(sl, rx, py + 1.0, rw, 4.0)
        p = para(tf, True)
        b = pitch.find('b')
        write(p, txt(b) + ' ', 1.15, INK, bold=True, line=1.4)
        write(p, txt(pitch).replace(txt(b), '', 1).strip(), 1.15, INK2, line=1.4)


# ================================================================ dispatch
def build(html='gtm-deck.html', out='PowerToFly-GTM-Sales-Deck.pptx'):
    soup = BeautifulSoup(open(os.path.join(HERE, html)), 'html.parser')
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.3333), Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = soup.select('.slide')
    for n, el in enumerate(slides, 1):
        sl = prs.slides.add_slide(blank)
        d = dark(el)
        bg = sl.background.fill
        bg.solid(); bg.fore_color.rgb = INK if d else WHITE
        if el.select_one('.hair'):
            hr = rect(sl, 0, 0 if 'top' in (el.select_one('.hair').get('style') or '')
                      else 100 * 9 / 16 - .5, 100, .5, GREEN_MID)

        has = lambda c: el.select_one('.' + c) is not None

        if has('covergrid'):                       L_cover(sl, el, d); continue
        if has('cta-center'):                      L_cta(sl, el, d);   continue
        if has('div-sub'):                         L_divider(sl, el, 0, d); continue
        if has('layers'):                          L_talent(sl, el, d); continue
        if has('soln'):                            L_dedicated(sl, el, d); continue

        y = chrome(sl, el, d)
        y = subline(sl, el, y, d)

        if   has('reachwrap'): L_reach(sl, el, y, d)
        elif has('stats'):     (L_stats5 if el.select_one('.stats.five') else L_stats)(sl, el, y, d)
        elif has('roles'):     L_roles(sl, el, y, d); L_facerow(sl, el, d)
        elif has('probs'):     L_probs(sl, el, y, d)
        elif has('svc'):       L_svc(sl, el, y, d)
        elif has('proc'):      L_proc(sl, el, y, d)
        elif has('bens'):      L_bens(sl, el, y, d)
        elif has('cmp'):       L_cmp(sl, el, y, d)
        elif has('audience'):  L_audience(sl, el, y, d)
        elif has('cal'):       L_cal(sl, el, y, d)
        elif has('media'):     L_media(sl, el, y, d)
        elif has('cases'):     L_cases(sl, el, y, d)
        elif has('tgrid'):     L_tgrid(sl, el, y, d)
        elif has('logos'):     L_logos(sl, el, y, d)
        elif has('pricing'):   L_pricing(sl, el, y, d)
        elif has('dtbl'):      L_dtbl(sl, el, y, d)
        elif has('grid6'):     L_grid6(sl, el, y, d)
        elif has('contacts'):  L_contacts(sl, el, y, d)
        elif has('twolist'):   L_twolist(sl, el, y, d)
        elif has('persona'):   L_persona(sl, el, y, d)
        else:
            print('  ! slide %d: no layout matched' % n)

    dest = os.path.join(HERE, out)
    prs.save(dest)
    problems = validate(dest)
    if problems:
        os.remove(dest)
        raise SystemExit('REFUSED TO SHIP — %d schema problems Google Slides would '
                         'reject:\n  ' % len(problems) + '\n  '.join(problems[:10]))
    return dest, len(slides)




# ================================================================ validation
def validate(path):
    """Refuse to ship a file Google Slides would reject.

    Google Slides validates the OOXML sequence strictly; PowerPoint does not.
    Every hand-written bit of XML in this file gets checked here, because a
    deck that only opens in PowerPoint is a deck that failed.
    """
    import zipfile
    from lxml import etree
    A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    SP = {'xfrm': 0, 'custGeom': 1, 'prstGeom': 1, 'noFill': 2, 'solidFill': 2,
          'gradFill': 2, 'blipFill': 2, 'pattFill': 2, 'grpFill': 2, 'ln': 3,
          'effectLst': 4, 'effectDag': 4, 'scene3d': 5, 'sp3d': 6, 'extLst': 9}
    TC = {'lnL': 0, 'lnR': 1, 'lnT': 2, 'lnB': 3, 'lnTlToBr': 4, 'lnBlToTr': 5,
          'cell3D': 6, 'noFill': 7, 'solidFill': 7, 'gradFill': 7, 'blipFill': 7,
          'pattFill': 7, 'grpFill': 7, 'headers': 8, 'extLst': 9}
    problems = []
    with zipfile.ZipFile(path) as z:
        for name in z.namelist():
            if not (name.startswith('ppt/slides/slide') and name.endswith('.xml')):
                continue
            root = etree.fromstring(z.read(name))
            for el in root.iter():
                local = etree.QName(el).localname
                if local not in ('spPr', 'tcPr'):
                    continue
                kids = [etree.QName(k).localname for k in el]
                table = SP if local == 'spPr' else TC
                if local == 'spPr' and kids.count('prstGeom') + kids.count('custGeom') > 1:
                    problems.append('%s: two geometries in one spPr (%s)' % (name, kids))
                seq = [table[k] for k in kids if k in table]
                if seq != sorted(seq):
                    problems.append('%s: %s out of schema order (%s)' % (name, local, kids))
    return problems


if __name__ == '__main__':
    path, n = build()
    print('wrote %s — %d slides, %.1f MB' % (os.path.basename(path), n,
                                             os.path.getsize(path) / 1e6))
